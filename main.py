import json
import os
import pathlib
import urllib
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from pptx.opc.packuri import PackURI
from pptx.parts.chart import ChartPart
from pptx.parts.embeddedpackage import EmbeddedXlsxPart
import copy
from copy import deepcopy
from datetime import datetime
from google.cloud import storage


def _get_blank_slide_layout(pres):
    layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
    min_items = min(layout_items_count)
    blank_layout_id = layout_items_count.index(min_items)
    return pres.slide_layouts[blank_layout_id]


def download_files(files, configuration_name):
    for index, val in enumerate(files):
        urllib.request.urlretrieve(val, "tmp/%s_%d.pptx" % (configuration_name, index))


def delete_files(files, configuration_name):
    for index, val in enumerate(files):
        os.remove("tmp/%s_%d.pptx" % (configuration_name, index))


def duplicate_slide_merge(pres, destination, index, isWordcloud):
    template = pres.slides[index]
    if isWordcloud:
        copied_slide = destination.slides.add_slide(destination.slide_layouts[9])
    else:
        copied_slide = destination.slides.add_slide(destination.slide_layouts[19])
    idx = 0
    imgDict = {}

    for shp in template.shapes:
        el = shp.element
        newel = deepcopy(el)
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            with open(str(len(destination.slides)) + '_' + shp.name + '.' + shp.image.ext, 'wb') as f:
                f.write(shp.image.blob)
            imgDict[str(len(destination.slides)) + '_' + shp.name + '.' + shp.image.ext] = [shp.left, shp.top,
                                                                                            shp.width, shp.height]
        elif shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            manageGroupShapes(shp, copied_slide, destination, imgDict, template)
        elif shp.shape_type == MSO_SHAPE_TYPE.CHART:
            newel.nvGraphicFramePr.cNvPr.id = 1000 + idx
            idx = idx + 1
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        elif shp.shape_type == MSO_SHAPE_TYPE.TABLE:
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        else:
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for key, value in template.part.rels.items():
        if "notesSlide`" not in value.reltype:
            target = value._target
            partnameArr = target.partname.split("/")
            partnameArr[-1] = 'chart_' + str(len(destination.slides)) + '_' + key + pathlib.Path(partnameArr[-1]).suffix
            partname = PackURI(str('/'.join(map(str, partnameArr))))
            target.partname = partname

            if "chart" in value.reltype:
                xlsx_blob = target.chart_workbook.xlsx_part.blob
                target = ChartPart(partname, target.content_type,
                                   copy.deepcopy(target._element), package=target.package)
                target.chart_workbook.xlsx_part = EmbeddedXlsxPart.new(
                    xlsx_blob, target.package)
                excelpartnameArr = target.chart_workbook.xlsx_part.partname.split("/")
                excelpartnameArr[-1] = 'Microsoft_Excel_Sheet_' + str(
                    len(destination.slides)) + '_' + key + pathlib.Path(
                    excelpartnameArr[-1]).suffix
                excelpartname = PackURI(str('/'.join(map(str, excelpartnameArr))))
                target.chart_workbook.xlsx_part.partname = excelpartname
                copied_slide.part.rels.add_relationship(value.reltype, target, value.rId)  # value.target_part
            else:
                if "image" not in value.reltype:
                    if not "xml" in str(value.target_ref):
                        if value.is_external:
                            copied_slide.part.rels.add_relationship(value.reltype, value.target_ref, value.rId,
                                                                    value.is_external)
                        else:
                            copied_slide.part.rels.add_relationship(value.reltype, value._target,
                                                                    value.rId)  # value.target_part
                else:
                    copied_slide.part.rels.add_relationship(value.reltype, target, value.rId)  # value.target_part

    # add pictures
    for k, v in imgDict.items():
        copied_slide.shapes.add_picture(k, v[0], v[1], v[2], v[3])
        os.remove(k)
    return copied_slide


def manageGroupShapes(shapes, copied_slide, destination, imgDict, parentShape):
    for shp in shapes.shapes:
        el = shp.element
        newel = deepcopy(el)
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            with open(str(len(destination.slides)) + '_' + str(shp.shape_id) + '_' + shp.name + '.' + shp.image.ext,
                      'wb') as f:
                f.write(shp.image.blob)
            newshape = copied_slide.shapes.add_picture(
                str(len(destination.slides)) + '_' + str(shp.shape_id) + '_' + shp.name + '.' + shp.image.ext, shp.left,
                shp.top, shp.width, shp.height)
            new_pic = newshape._element
            new_pic.getparent().remove(new_pic)
            os.remove(str(len(destination.slides)) + '_' + str(shp.shape_id) + '_' + shp.name + '.' + shp.image.ext)

        elif shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            manageGroupShapes(shp, copied_slide, destination, imgDict, shapes)
    return copied_slide


def pptxmerge(merge_files, out_file):
    filename = "survey_template"
    download_files(merge_files, filename)
    for index, file in enumerate(merge_files):
        if index == 0:
            destination = Presentation("tmp/" + filename + "_%d.pptx" % (index))
        else:
            source = Presentation("tmp/" + filename + "_%d.pptx" % (index))
            for i in range(len(source.slides)):
                if "word_cloud" in merge_files[index]:
                    duplicate_slide_merge(source, destination, i, True)
                else:
                    duplicate_slide_merge(source, destination, i, False)
    destination.save(out_file)
    delete_files(merge_files, filename)
    return out_file


if __name__ == '__main__':
    # request = '{"configurationName":"SV_6PCcJpGbtHreLHL_bhavin_test","data":[{"categoryData":["In-person sales representative","Phone","Live chat","In- storerepresentative","Email","Website","Information or signage in store"],"multiple_segments":[{"noOfUsers":11,"seriesData":[[0.4545,0.0909,0.2727,0.0909,0.0909,0,0],[0.0909,0.1818,0.1818,0.1818,0.0909,0,0],[0,0.2727,0.0909,0.0909,0.0909,0,0],[0.0909,0.0909,0.0909,0,0.0909,0,0],[0,0,0,0,0,0.1818,0]]},{"noOfUsers":9,"seriesData":[[0.5556,0.1111,0.1111,0.1111,0.1111,0,0],[0.1111,0.1111,0.2222,0.2222,0,0,0],[0,0.2222,0.1111,0.1111,0.1111,0,0],[0.1111,0.1111,0.1111,0,0.1111,0,0],[0,0,0,0,0,0.2222,0]]}],"imagesData":[],"topRankValue":5,"noOfChoicePerSlide":5,"questionText":"QID201 (Q130)-What type(s) of contact were most helpful?","questionAnalysisType":"Rank Interface","segmentName":["filter1","filter4"],"segmentColor":["23735D","751541","203673","568D11","0C779D"]}]}'
    # data = json.loads(request)
    result = {}
    # ppts_files = ['tmp/1.pptx', 'tmp/2.pptx','tmp/3.pptx','tmp/4.pptx','tmp/5.pptx','tmp/6.pptx','tmp/7.pptx','tmp/8.pptx'];
    ppts_files = [
        "https://storage.googleapis.com/intuify_production/analysis/ppt/path_to_purchase/SV_6PCcJpGbtHreLHL_bhavin_test/SV_6PCcJpGbtHreLHL_bhavin_test_01_11_2021_10_29_13.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/word_cloud/SV_6PCcJpGbtHreLHL_ExampleWordCloud/SV_6PCcJpGbtHreLHL_ExampleWordCloud_01_11_2021_10_31_25.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/sequential_select/SV_79G6akh0lWqqMPs_sequential_checkup/SV_79G6akh0lWqqMPs_sequential_checkup_06_12_2021_06_35_54.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/rank_interface/SV_6PCcJpGbtHreLHL_FrontierFiber/SV_6PCcJpGbtHreLHL_FrontierFiber_01_11_2021_10_37_18.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/sentiment_emotions/SV_6PCcJpGbtHreLHL_630/SV_6PCcJpGbtHreLHL_630_01_11_2021_10_33_32.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/three_way_swipe_interface/SV_6PCcJpGbtHreLHL_bhavin_test1/SV_6PCcJpGbtHreLHL_bhavin_test1_02_11_2021_05_27_38.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/single_select_multi_choice/SV_6PCcJpGbtHreLHL_bhavin_test/SV_6PCcJpGbtHreLHL_bhavin_test_01_11_2021_10_44_21.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/time_based_rating/SV_8crCfzUH6DHY4x8_durgesh_test11/SV_8crCfzUH6DHY4x8_durgesh_test11_03_12_2021_10_46_34.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/agreement_battery/SV_6PCcJpGbtHreLHL_meet_test/SV_6PCcJpGbtHreLHL_meet_test_06_12_2021_06_19_50.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/time_based_rating/SV_6PCcJpGbtHreLHL_meet_test_withoutsegment_bidirection/SV_6PCcJpGbtHreLHL_meet_test_withoutsegment_bidirection_07_12_2021_13_20_39.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/tap_and_talk_v3/SV_5vj6JY1uxPs7p0p_Test/SV_5vj6JY1uxPs7p0p_Test_02_08_2021_13_06_43.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/multi_swipe_bar_stack/SV_79G6akh0lWqqMPs_temp/SV_79G6akh0lWqqMPs_temp_07_12_2021_13_24_46.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/multi_swipe_heat_map/SV_79G6akh0lWqqMPs_bhavin_test/SV_79G6akh0lWqqMPs_bhavin_test_07_12_2021_13_25_41.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/grid_statement/SV_79G6akh0lWqqMPs_grid_checkup/SV_79G6akh0lWqqMPs_grid_checkup_06_12_2021_06_38_31.pptx",
        "https://storage.googleapis.com/intuify_production/analysis/ppt/time_based_rating/SV_blLWfIpDiq6LTLf_meet_test/SV_blLWfIpDiq6LTLf_meet_test_07_12_2021_13_28_53.pptx"
    ]
    now = datetime.now()
    result['data'] = {
        'url': pptxmerge(ppts_files, 'tmp/' + now.strftime("%d_%m_%Y_%H_%M_%S") + '.pptx')
    }
    print(result)

    # return json.dumps(result)
